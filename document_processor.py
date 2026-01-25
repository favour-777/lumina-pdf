# Document Processor - Handles downloading and text extraction from various formats

import aiohttp
import hashlib
from pathlib import Path
from io import BytesIO
import PyPDF2
import docx
import mammoth
from bs4 import BeautifulSoup
import ebooklib
from ebooklib import epub
from striprtf.striprtf import rtf_to_text
from pptx import Presentation
from openpyxl import load_workbook
import xlrd
import re


class DocumentProcessor:
    """Extract text from various document formats"""
    
    SUPPORTED_FORMATS = {
        'pdf': ['.pdf'],
        'docx': ['.docx', '.doc'],
        'pptx': ['.pptx', '.ppt'],
        'xlsx': ['.xlsx'],
        'xls': ['.xls'],
        'txt': ['.txt', '.text'],
        'markdown': ['.md', '.markdown'],
        'html': ['.html', '.htm'],
        'epub': ['.epub'],
        'rtf': ['.rtf']
    }
    
    def __init__(self):
        self.session = None
    
    async def download_file(self, url: str) -> dict:
        """Download file from URL and detect format"""
        if not self.session:
            timeout = aiohttp.ClientTimeout(total=300)
            self.session = aiohttp.ClientSession(timeout=timeout)
        
        async with self.session.get(url) as response:
            if response.status != 200:
                raise Exception(f"Failed to download file: HTTP {response.status}")
            
            content = await response.read()
            
            filename = self._extract_filename(url, response.headers)
            file_format = self._detect_format(filename, content)
            
            file_id = hashlib.md5(content).hexdigest()[:12]
            
            return {
                'content': content,
                'filename': filename,
                'format': file_format,
                'file_id': file_id,
                'size': len(content)
            }
    
    def _extract_filename(self, url: str, headers: dict) -> str:
        """Extract filename from URL or headers"""
        if 'Content-Disposition' in headers:
            content_disp = headers['Content-Disposition']
            if 'filename=' in content_disp:
                filename = content_disp.split('filename=')[-1].strip('"\'')
                return filename
        
        from urllib.parse import urlparse, unquote
        path = urlparse(url).path
        filename = Path(unquote(path)).name
        
        return filename or 'document'
    
    def _detect_format(self, filename: str, content: bytes) -> str:
        """Detect document format from filename or content"""
        ext = Path(filename).suffix.lower()
        
        for format_type, extensions in self.SUPPORTED_FORMATS.items():
            if ext in extensions:
                return format_type
        
        if content.startswith(b'%PDF'):
            return 'pdf'
        elif content.startswith(b'PK\x03\x04'):
            if b'word/' in content[:1000]:
                return 'docx'
            elif b'ppt/' in content[:1000]:
                return 'pptx'
            elif b'xl/' in content[:1000]:
                return 'xlsx'
            elif b'EPUB' in content[:100]:
                return 'epub'
        elif content.startswith(b'\xd0\xcf\x11\xe0'):
            if ext == '.xls':
                return 'xls'
            elif ext == '.doc':
                return 'docx'
            elif ext == '.ppt':
                return 'pptx'
            return 'xls'
        elif content.startswith(b'{\\rtf'):
            return 'rtf'
        elif b'<html' in content[:1000].lower():
            return 'html'
        
        return 'txt'
    
    async def extract_text(self, content: bytes, format_type: str) -> str:
        """Extract text from document based on format"""
        extractors = {
            'pdf': self._extract_pdf,
            'docx': self._extract_docx,
            'pptx': self._extract_pptx,
            'xlsx': self._extract_xlsx,
            'xls': self._extract_xls,
            'txt': self._extract_txt,
            'markdown': self._extract_txt,
            'html': self._extract_html,
            'epub': self._extract_epub,
            'rtf': self._extract_rtf
        }
        
        extractor = extractors.get(format_type, self._extract_txt)
        text = await extractor(content)
        
        text = self._clean_text(text)
        
        return text
    
    async def _extract_pdf(self, content: bytes) -> str:
        """Extract text from PDF"""
        pdf_file = BytesIO(content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        text_parts = []
        for page in pdf_reader.pages:
            text_parts.append(page.extract_text())
        
        return '\n\n'.join(text_parts)
    
    async def _extract_docx(self, content: bytes) -> str:
        """Extract text from DOCX"""
        docx_file = BytesIO(content)
        
        result = mammoth.convert_to_html(docx_file)
        html_text = result.value
        
        soup = BeautifulSoup(html_text, 'html.parser')
        text = soup.get_text(separator='\n')
        
        return text
    
    async def _extract_txt(self, content: bytes) -> str:
        """Extract text from plain text files"""
        for encoding in ['utf-8', 'latin-1', 'cp1252']:
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        
        return content.decode('utf-8', errors='ignore')
    
    async def _extract_html(self, content: bytes) -> str:
        """Extract text from HTML"""
        html = await self._extract_txt(content)
        soup = BeautifulSoup(html, 'html.parser')
        
        for script in soup(['script', 'style']):
            script.decompose()
        
        text = soup.get_text(separator='\n')
        return text
    
    async def _extract_epub(self, content: bytes) -> str:
        """Extract text from EPUB"""
        epub_file = BytesIO(content)
        book = epub.read_epub(epub_file)
        
        text_parts = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                text_parts.append(soup.get_text())
        
        return '\n\n'.join(text_parts)
    
    async def _extract_rtf(self, content: bytes) -> str:
        """Extract text from RTF"""
        rtf_text = await self._extract_txt(content)
        text = rtf_to_text(rtf_text)
        return text
    
    async def _extract_pptx(self, content: bytes) -> str:
        """Extract text from PowerPoint"""
        pptx_file = BytesIO(content)
        prs = Presentation(pptx_file)
        
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n--- Slide {slide_num} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        
        return '\n\n'.join(text_parts)
    
    async def _extract_xlsx(self, content: bytes) -> str:
        """Extract text from Excel (new format)"""
        xlsx_file = BytesIO(content)
        wb = load_workbook(xlsx_file, data_only=True)
        
        text_parts = []
        for sheet in wb.worksheets:
            text_parts.append(f"\n=== Sheet: {sheet.title} ===\n")
            for row in sheet.iter_rows(values_only=True):
                row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                if row_text.strip():
                    text_parts.append(row_text)
        
        return '\n'.join(text_parts)
    
    async def _extract_xls(self, content: bytes) -> str:
        """Extract text from Excel (old .xls format)"""
        xls_file = BytesIO(content)
        wb = xlrd.open_workbook(file_contents=content)
        
        text_parts = []
        for sheet in wb.sheets():
            text_parts.append(f"\n=== Sheet: {sheet.name} ===\n")
            for row_idx in range(sheet.nrows):
                row = sheet.row_values(row_idx)
                row_text = '\t'.join([str(cell) if cell else '' for cell in row])
                if row_text.strip():
                    text_parts.append(row_text)
        
        return '\n'.join(text_parts)
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize extracted text"""
        text = re.sub(r'\n\s*\n', '\n\n', text)
        text = re.sub(r'[ \t]+', ' ', text)
        text = re.sub(r'\n\s*\d+\s*\n', '\n', text)
        text = text.replace('\r\n', '\n')
        
        return text.strip()
    
    async def close(self):
        """Close aiohttp session"""
        if self.session:
            await self.session.close()